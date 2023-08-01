import io
import os
import re
from pptx import Presentation
from pptx.util import Inches
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build

class GoogleSlidesPresentation:
    def __init__(self, credentials_path=r"C:\Users\porwa\Downloads\client_secret_524550495379-ijrb0lbfal2j8cel6bjii6or8qq3j5ta.apps.googleusercontent.com.json"):
        self.credentials_path = credentials_path

    def authenticate(self):
        SCOPES = ['https://www.googleapis.com/auth/presentations']

        creds = None
        if os.path.exists('token.json'):
            creds = Credentials.from_authorized_user_file('token.json', SCOPES)

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                flow = InstalledAppFlow.from_client_secrets_file(
                    self.credentials_path, SCOPES)
                creds = flow.run_local_server(port=0)
            with open('token.json', 'w') as token:
                token.write(creds.to_json())

        return creds

class SlideGenerator:
    def __init__(self, content, max_words_per_slide=200):
        self.content = content
        self.max_words_per_slide = max_words_per_slide
        self.pages = self.split_text_into_pages(content, max_words_per_slide)

    def split_text_into_pages(self, text, max_words_per_page):
        # Split the text into chunks based on word limit
        words = text.split()
        pages = []
        current_page = ""
        for word in words:
            if len(current_page.split()) + 1 <= max_words_per_page:
                current_page += word + " "
            else:
                pages.append(current_page)
                current_page = word + " "
        if current_page:
            pages.append(current_page)
        return pages

    def create_presentation(self):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Title and Content layout

        for page in self.pages:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            # Set title and content
            title.text = "Page {}".format(self.pages.index(page) + 1)
            content.text = page

        return prs

class GoogleSlidesUploader:
    def __init__(self, credentials_path='path/to/your/credentials.json'):
        self.presentation = None
        self.slides_service = None
        self.credentials_path = credentials_path

    def authenticate_google_slides(self):
        creds = GoogleSlidesPresentation(self.credentials_path).authenticate()
        self.slides_service = build('slides', 'v1', credentials=creds)

    def save_to_google_slides(self, presentation):
        self.authenticate_google_slides()
        presentation_request = self.slides_service.presentations().create(body={"title": "My Presentation"})
        presentation_response = presentation_request.execute()
        presentation_id = presentation_response['presentationId']

        pptx_bytes = io.BytesIO()
        presentation.save(pptx_bytes)
        pptx_bytes.seek(0)

        self.slides_service.presentations().batchUpdate(
            presentationId=presentation_id,
            body={
                "requests": [
                    {
                        "createSlide": {
                            "insertionIndex": 1,
                            "slideLayout": {
                                "predefinedLayout": "BLANK"
                            }
                        }
                    }
                ]
            }
        ).execute()

        self.slides_service.presentations().pages().upload(
            presentationId=presentation_id,
            body={},
            media_body=io.BytesIO(pptx_bytes.read()),
        ).execute()

if __name__ == "__main__":
    # Define your content here (replace with your own text)
    content = """
    This is the first page with some content. It should not exceed 150-200 words.
    You can add more paragraphs as needed. Make sure the text is concise and relevant.

    This is the second page with some more content. Keep the information clear and to the point.
    Remember to break the content into smaller paragraphs to fit within the word limit.

    And so on for the rest of the content...
    """

    slide_generator = SlideGenerator(content)
    presentation = slide_generator.create_presentation()

    uploader = GoogleSlidesUploader()
    uploader.save_to_google_slides(presentation)
